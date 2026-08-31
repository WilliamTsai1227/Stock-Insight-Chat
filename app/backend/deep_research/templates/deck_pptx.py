"""
簡報的 PowerPoint 產出（.pptx）。

與 HTML 版吃同一份 `DeckDoc`，一頁對應一張原生投影片，`note` 進講者備忘稿。

版面全部用空白版面配置 + 自己擺文字框畫出來，不套 PowerPoint 內建版面 ——
內建版面的佔位框會帶進範本自己的字型與色彩，主題（themes.py）就管不到了。
自己畫的代價是要處理座標，換到的是「HTML 與 pptx 看起來是同一份簡報」。

視覺會比 HTML 版樸素：CSS 的漸層光暈、翻頁動畫這些換不過來，也不該換。
pptx 的價值是能在會議室用 PowerPoint 開、能直接改、講稿在該在的地方。
"""

from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Any, List, Sequence, Tuple

from .office import hex_rgb, office_font
from .themes import Theme

# 16:9，13.333 x 7.5 英寸（PowerPoint 的寬螢幕預設）
_SLIDE_W = 13.333
_SLIDE_H = 7.5
_MARGIN = 0.85
_CONTENT_W = _SLIDE_W - _MARGIN * 2


def _rgb(color: str):
    from pptx.dml.color import RGBColor

    return RGBColor(*hex_rgb(color))


def _set_font(run, latin: str, east_asian: str) -> None:
    """
    同時設定 latin 與 eastAsia 字型。

    `run.font.name` 只寫 a:latin；中文字元走 a:ea，沒設的話 PowerPoint 會自己
    挑字型，同一頁的中英文會對不齊基線。
    """
    from pptx.oxml.ns import qn

    run.font.name = latin
    rpr = run.font._rPr
    for tag in ("a:ea", "a:cs"):
        existing = rpr.find(qn(tag))
        if existing is None:
            existing = rpr.makeelement(qn(tag), {})
            rpr.append(existing)
        existing.set("typeface", east_asian if tag == "a:ea" else latin)


class _Deck:
    """一份簡報的畫布：所有座標與樣式決策集中在這裡。"""

    def __init__(self, presentation, theme: Theme, deck_title: str, total: int) -> None:
        from pptx.util import Inches, Pt

        self.presentation = presentation
        self.theme = theme
        self.palette = theme.tokens
        self.display = office_font(theme, "display")
        self.body = office_font(theme, "body")
        self.deck_title = deck_title
        self.total = total
        self.Inches = Inches
        self.Pt = Pt
        self.blank = presentation.slide_layouts[6]

    # ── 基礎 ────────────────────────────────────────────────
    def new_slide(self, note: str):
        slide = self.presentation.slides.add_slide(self.blank)
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = _rgb(self.palette["bg"])
        if note and note.strip():
            slide.notes_slide.notes_text_frame.text = note.strip()
        return slide

    def textbox(self, slide, left: float, top: float, width: float, height: float):
        from pptx.util import Inches

        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        return frame

    def write(
        self,
        frame,
        text: str,
        *,
        size: float,
        color: str,
        bold: bool = False,
        fonts: Tuple[str, str] | None = None,
        space_after: float = 0,
        line_spacing: float = 1.25,
        first: bool = False,
        align: str | None = None,
    ):
        """在文字框裡加一段。`first=True` 用掉文字框自帶的空段落，不留空行。"""
        from pptx.enum.text import PP_ALIGN

        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        paragraph.line_spacing = line_spacing
        paragraph.space_after = self.Pt(space_after)
        if align:
            paragraph.alignment = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT}[align]

        run = paragraph.add_run()
        run.text = text
        run.font.size = self.Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        _set_font(run, *(fonts or self.body))
        return paragraph

    def rect(self, slide, left: float, top: float, width: float, height: float, color: str,
             *, rounded: bool = False, outline: str | None = None):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(color)
        if outline:
            shape.line.color.rgb = _rgb(outline)
            shape.line.width = self.Pt(0.75)
        else:
            shape.line.fill.background()
        # 預設 autoshape 會帶範本的陰影；主題裡沒有陰影這個維度，關掉才一致
        shape.shadow.inherit = False
        if rounded:
            shape.adjustments[0] = 0.08
        return shape

    # ── 共用元件 ────────────────────────────────────────────
    def kicker(self, slide, text: str, top: float = 0.62) -> None:
        if not text:
            return
        frame = self.textbox(slide, _MARGIN, top, _CONTENT_W, 0.3)
        self.write(
            frame,
            text,
            size=11,
            color=self.palette["accent"],
            bold=True,
            first=True,
            line_spacing=1.0,
        )

    def rule(self, slide, top: float, width: float = 1.5) -> None:
        self.rect(slide, _MARGIN, top, width, 0.045, self.palette["accent"])

    def foot(self, slide, index: int) -> None:
        """頁尾：簡報名 + 頁碼 + 底部進度條（對應 HTML 版的 .slide-foot）。"""
        frame = self.textbox(slide, _MARGIN, _SLIDE_H - 0.62, _CONTENT_W - 1.2, 0.28)
        self.write(
            frame,
            self.deck_title,
            size=9,
            color=self.palette["muted"],
            first=True,
            line_spacing=1.0,
        )

        number = self.textbox(slide, _SLIDE_W - _MARGIN - 1.2, _SLIDE_H - 0.62, 1.2, 0.28)
        self.write(
            number,
            f"{index} / {self.total}",
            size=9,
            color=self.palette["muted"],
            first=True,
            line_spacing=1.0,
            align="center",
        )

        self.rect(slide, 0, _SLIDE_H - 0.06, _SLIDE_W, 0.06, self.palette["line"])
        self.rect(
            slide, 0, _SLIDE_H - 0.06, _SLIDE_W * index / self.total, 0.06,
            self.palette["accent"],
        )


# ─────────────────────────────────────────────────────────────
# 版面
# ─────────────────────────────────────────────────────────────
def _bullet_size(bullets: Sequence[str]) -> Tuple[float, float]:
    """
    依條列數量與長度決定字級與行距。

    與 HTML 版的 `_density()` 同一個取向：不精確量測，保守估計，
    寧可字小一號也不要溢出畫布 —— pptx 溢出的字是直接被裁掉，救不回來。
    """
    count = len(bullets)
    chars = sum(len(b) for b in bullets)
    if count >= 5 or chars > 320:
        return 15, 0.30
    if count >= 4 or chars > 220:
        return 17, 0.34
    return 19, 0.38


def _title_size(title: str) -> float:
    if len(title) > 34:
        return 26
    if len(title) > 22:
        return 30
    return 34


def _render_cover(deck: _Deck, slide, item: Any) -> None:
    deck.rect(slide, 0, 0, 0.18, _SLIDE_H, deck.palette["accent"])

    deck.kicker(slide, item.kicker or "深度研究", top=1.9)
    frame = deck.textbox(slide, _MARGIN, 2.35, _CONTENT_W - 1.5, 2.4)
    deck.write(
        frame,
        item.title,
        size=42 if len(item.title) <= 28 else 34,
        color=deck.palette["ink"],
        bold=True,
        fonts=deck.display,
        first=True,
        line_spacing=1.12,
    )

    if item.lead:
        lead = deck.textbox(slide, _MARGIN, 4.95, _CONTENT_W - 2.5, 1.0)
        deck.write(
            lead, item.lead, size=17, color=deck.palette["ink2"], first=True, line_spacing=1.4
        )

    stamp = deck.textbox(slide, _MARGIN, _SLIDE_H - 1.0, _CONTENT_W, 0.3)
    deck.write(
        stamp,
        f"{datetime.now().strftime('%Y-%m-%d')}　·　Insight 深度研究",
        size=10,
        color=deck.palette["muted"],
        first=True,
        line_spacing=1.0,
    )


def _render_section(deck: _Deck, slide, item: Any, index: int, section_no: int) -> None:
    number = deck.textbox(slide, _MARGIN, 2.5, 2.2, 1.4)
    deck.write(
        number,
        f"{section_no:02d}",
        size=64,
        color=deck.palette["accent"],
        bold=True,
        fonts=deck.display,
        first=True,
        line_spacing=1.0,
    )

    frame = deck.textbox(slide, _MARGIN, 3.85, _CONTENT_W - 1.5, 1.6)
    deck.write(
        frame,
        item.title,
        size=_title_size(item.title),
        color=deck.palette["ink"],
        bold=True,
        fonts=deck.display,
        first=True,
        line_spacing=1.18,
    )
    if item.lead:
        deck.write(
            frame, item.lead, size=15, color=deck.palette["body"], space_after=0,
            line_spacing=1.4,
        )
    deck.foot(slide, index)


def _render_quote(deck: _Deck, slide, item: Any, index: int) -> None:
    mark = deck.textbox(slide, _MARGIN, 1.55, 1.2, 1.2)
    deck.write(
        mark, "“", size=80, color=deck.palette["accent"], bold=True,
        fonts=deck.display, first=True, line_spacing=1.0,
    )

    text = item.lead or item.title
    frame = deck.textbox(slide, _MARGIN, 2.75, _CONTENT_W - 1.2, 2.2)
    deck.write(
        frame,
        text,
        size=30 if len(text) <= 40 else 24,
        color=deck.palette["ink"],
        bold=True,
        fonts=deck.display,
        first=True,
        line_spacing=1.35,
    )

    attribution = deck.textbox(slide, _MARGIN, 5.35, _CONTENT_W - 1.2, 0.4)
    deck.write(
        attribution,
        item.kicker or "本次研究的核心結論",
        size=12,
        color=deck.palette["muted"],
        first=True,
        line_spacing=1.0,
    )
    deck.foot(slide, index)


def _render_stats(deck: _Deck, slide, item: Any, index: int) -> None:
    _render_head(deck, slide, item)

    stats = list(item.stats)[:4]
    if not stats:
        return
    gap = 0.3
    width = (_CONTENT_W - gap * (len(stats) - 1)) / len(stats)
    rounded = deck.palette.get("radius", "8px") != "0px"

    for n, stat in enumerate(stats):
        left = _MARGIN + n * (width + gap)
        deck.rect(
            slide, left, 3.5, width, 1.95, deck.palette["accent-soft"],
            rounded=rounded, outline=deck.palette["accent-line"],
        )
        frame = deck.textbox(slide, left + 0.32, 3.85, width - 0.64, 1.3)
        deck.write(
            frame,
            stat.value,
            size=34 if len(stat.value) <= 8 else 22,
            color=deck.palette["accent"],
            bold=True,
            fonts=deck.display,
            first=True,
            line_spacing=1.05,
            space_after=6,
        )
        deck.write(
            frame, stat.label, size=12, color=deck.palette["body"], line_spacing=1.3
        )
    deck.foot(slide, index)


def _render_compare(deck: _Deck, slide, item: Any, index: int) -> None:
    _render_head(deck, slide, item)

    columns = list(item.columns)[:2]
    if not columns:
        return
    gap = 0.55
    width = (_CONTENT_W - gap) / 2

    for n, column in enumerate(columns):
        left = _MARGIN + n * (width + gap)
        deck.rect(slide, left, 3.45, width, 0.045, deck.palette["accent"])
        frame = deck.textbox(slide, left, 3.7, width, 2.6)
        deck.write(
            frame,
            column.heading,
            size=18,
            color=deck.palette["ink"],
            bold=True,
            fonts=deck.display,
            first=True,
            line_spacing=1.2,
            space_after=10,
        )
        size, space = _bullet_size(column.bullets)
        for bullet in column.bullets:
            deck.write(
                frame,
                f"·　{bullet}",
                size=min(size, 15),
                color=deck.palette["body"],
                space_after=space * 24,
                line_spacing=1.35,
            )
    deck.foot(slide, index)


def _render_agenda(deck: _Deck, slide, item: Any, index: int) -> None:
    _render_head(deck, slide, item)

    frame = deck.textbox(slide, _MARGIN, 3.45, _CONTENT_W - 1.0, 3.2)
    size, space = _bullet_size(item.bullets)
    for n, bullet in enumerate(item.bullets, start=1):
        deck.write(
            frame,
            f"{n:02d}　{bullet}",
            size=size,
            color=deck.palette["ink2"],
            space_after=space * 30,
            line_spacing=1.35,
            first=(n == 1),
        )
    deck.foot(slide, index)


def _render_head(deck: _Deck, slide, item: Any) -> None:
    """一般內容頁的頁首：小標籤、頁標題、補充句、色帶。"""
    deck.kicker(slide, item.kicker or "重點")

    frame = deck.textbox(slide, _MARGIN, 1.05, _CONTENT_W - 1.0, 1.5)
    deck.write(
        frame,
        item.title,
        size=_title_size(item.title),
        color=deck.palette["ink"],
        bold=True,
        fonts=deck.display,
        first=True,
        line_spacing=1.16,
    )

    top = 2.55
    if item.lead and item.lead.strip():
        lead = deck.textbox(slide, _MARGIN, 2.45, _CONTENT_W - 1.6, 0.6)
        deck.write(
            lead, item.lead, size=14, color=deck.palette["body"], first=True, line_spacing=1.4
        )
        top = 3.1
    deck.rule(slide, top)


def _render_bullets(deck: _Deck, slide, item: Any, index: int) -> None:
    _render_head(deck, slide, item)

    if item.bullets:
        frame = deck.textbox(slide, _MARGIN, 3.5, _CONTENT_W - 0.8, 3.0)
        size, space = _bullet_size(item.bullets)
        for n, bullet in enumerate(item.bullets):
            deck.write(
                frame,
                f"·　{bullet}",
                size=size,
                color=deck.palette["ink2"],
                space_after=space * 30,
                line_spacing=1.38,
                first=(n == 0),
            )
    deck.foot(slide, index)


def _render_closing(deck: _Deck, slide, item: Any, index: int) -> None:
    deck.kicker(slide, item.kicker or "結語", top=1.75)

    frame = deck.textbox(slide, _MARGIN, 2.15, _CONTENT_W - 1.2, 1.5)
    deck.write(
        frame,
        item.title,
        size=_title_size(item.title),
        color=deck.palette["ink"],
        bold=True,
        fonts=deck.display,
        first=True,
        line_spacing=1.16,
    )
    if item.lead:
        deck.write(
            frame, item.lead, size=15, color=deck.palette["body"], line_spacing=1.4
        )

    if item.bullets:
        body = deck.textbox(slide, _MARGIN, 4.1, _CONTENT_W - 1.0, 2.3)
        size, space = _bullet_size(item.bullets)
        for n, bullet in enumerate(item.bullets):
            deck.write(
                body,
                f"·　{bullet}",
                size=min(size, 16),
                color=deck.palette["ink2"],
                space_after=space * 26,
                line_spacing=1.35,
                first=(n == 0),
            )
    deck.foot(slide, index)


def render_deck_pptx(doc: Any, *, model: str, query: str, theme: Theme) -> bytes:
    """把 `DeckDoc` 轉成 .pptx 的位元組內容。"""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(_SLIDE_W)
    presentation.slide_height = Inches(_SLIDE_H)

    core = presentation.core_properties
    core.title = doc.title
    core.subject = query
    core.comments = f"由 Insight 深度研究以 {model} 產生"

    deck = _Deck(presentation, theme, doc.title, max(len(doc.slides), 1))
    section_no = 0

    for index, item in enumerate(doc.slides, start=1):
        slide = deck.new_slide(item.note)
        layout = item.layout

        if layout == "cover":
            _render_cover(deck, slide, item)
        elif layout == "section":
            section_no += 1
            _render_section(deck, slide, item, index, section_no)
        elif layout == "quote":
            _render_quote(deck, slide, item, index)
        elif layout == "stats" and item.stats:
            _render_stats(deck, slide, item, index)
        elif layout == "compare" and item.columns:
            _render_compare(deck, slide, item, index)
        elif layout == "agenda" and item.bullets:
            _render_agenda(deck, slide, item, index)
        elif layout == "closing":
            _render_closing(deck, slide, item, index)
        else:
            # stats / compare / agenda 缺對應欄位時也走這裡，不會生出一頁空白
            _render_bullets(deck, slide, item, index)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
